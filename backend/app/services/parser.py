"""Document parsing service - inspired by MinerU's approach."""

import os
from typing import List, Dict
from pathlib import Path


class DocumentParser:
    """Parse documents into structured text with paragraph preservation."""

    @staticmethod
    async def parse(file_path: str, file_type: str) -> str:
        """Parse a document and return its full text content."""
        parsers = {
            "pdf": DocumentParser._parse_pdf,
            "docx": DocumentParser._parse_docx,
            "txt": DocumentParser._parse_txt,
            "md": DocumentParser._parse_txt,
            "pptx": DocumentParser._parse_pptx,
            "xlsx": DocumentParser._parse_xlsx,
            "csv": DocumentParser._parse_csv,
            "image": DocumentParser._parse_image,
            "audio": DocumentParser._parse_media,
            "video": DocumentParser._parse_media,
        }
        parser = parsers.get(file_type)
        if not parser:
            raise ValueError(f"Unsupported file type: {file_type}")

        text = await parser(file_path)
        return text.strip()

    @staticmethod
    async def _parse_pdf(file_path: str) -> str:
        """Parse PDF using PyMuPDF, preserving paragraph structure."""
        import fitz  # PyMuPDF

        doc = fitz.open(file_path)
        full_text = []
        for page_num, page in enumerate(doc):
            # Get text blocks to preserve paragraph structure
            blocks = page.get_text("blocks")
            page_text = []
            for block in blocks:
                text = block[4].strip() if len(block) > 4 else ""
                if text:
                    page_text.append(text)
            if page_text:
                full_text.append("\n\n".join(page_text))
        doc.close()
        return "\n\n".join(full_text)

    @staticmethod
    async def _parse_docx(file_path: str) -> str:
        """Parse DOCX file."""
        from docx import Document

        doc = Document(file_path)
        paragraphs = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                paragraphs.append(text)

        # Also extract text from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    if cell.text.strip():
                        row_text.append(cell.text.strip())
                if row_text:
                    paragraphs.append(" | ".join(row_text))

        return "\n\n".join(paragraphs)

    @staticmethod
    async def _parse_txt(file_path: str) -> str:
        """Parse plain text / markdown files."""
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()

    @staticmethod
    async def _parse_pptx(file_path: str) -> str:
        """Parse PPTX file."""
        from pptx import Presentation

        prs = Presentation(file_path)
        slides_text = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_text.append(text)
            if slide_text:
                slides_text.append(f"--- Slide {slide_num} ---\n" + "\n".join(slide_text))
        return "\n\n".join(slides_text)

    @staticmethod
    async def _parse_xlsx(file_path: str) -> str:
        """Parse Excel file, extracting all sheets."""
        from openpyxl import load_workbook
        wb = load_workbook(file_path, read_only=True, data_only=True)
        parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                rows.append("\t".join(str(c or "") for c in row))
            parts.append(f"## Sheet: {sheet_name}\n" + "\n".join(rows))
        wb.close()
        return "\n\n".join(parts)

    @staticmethod
    async def _parse_csv(file_path: str) -> str:
        """Parse CSV file."""
        import csv
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            reader = csv.reader(f)
            rows = ["\t".join(row) for row in reader]
        return "\n".join(rows)

    @staticmethod
    async def _parse_media(file_path: str) -> str:
        """Parse media file - returns metadata placeholder."""
        import os
        filename = os.path.basename(file_path)
        file_size = os.path.getsize(file_path)
        return (
            f"[媒体文件]\n文件名: {filename}\n大小: {file_size} bytes\n"
            f"说明: 音频/视频文件已上传。语音转文字将在后续版本支持。"
        )

    @staticmethod
    async def _parse_image(file_path: str) -> str:
        """Parse image file - returns basic metadata + OCR placeholder."""
        import os
        filename = os.path.basename(file_path)
        file_size = os.path.getsize(file_path)
        return (
            f"[图片文件]\n"
            f"文件名: {filename}\n"
            f"大小: {file_size} bytes\n"
            f"说明: 图片已上传，OCR文字识别将在后续版本支持。"
            f"如需检索此图片内容，请在FAQ中添加相关问答对，或为此图片添加文字描述。"
        )

    @staticmethod
    async def _extract_wechat_article(soup: "BeautifulSoup") -> str | None:
        """Extract content from mp.weixin.qq.com articles.

        WeChat stores the article body as an escaped HTML string inside a
        ``<script>`` tag (``cgiDataNew.content_noencode``).  The visible DOM
        is rendered client-side, so BeautifulSoup never sees it.
        """
        import re
        import html as html_mod

        def _decode_js_escapes(raw: str) -> str:
            """Decode JS escape sequences character-by-character.

            Only converts backslash-escapes (\\xHH, \\uHHHH, \\n, \\t).
            Literal Chinese characters pass through as-is — no re-encoding.
            """
            parts: list[str] = []
            i = 0
            n = len(raw)
            while i < n:
                ch = raw[i]
                if ch == "\\" and i + 1 < n:
                    nxt = raw[i + 1]
                    if nxt == "x" and i + 3 < n:
                        hex_s = raw[i + 2 : i + 4]
                        try:
                            parts.append(chr(int(hex_s, 16)))
                            i += 4
                            continue
                        except ValueError:
                            pass
                    elif nxt == "u" and i + 5 < n:
                        hex_s = raw[i + 2 : i + 6]
                        try:
                            parts.append(chr(int(hex_s, 16)))
                            i += 6
                            continue
                        except ValueError:
                            pass
                    elif nxt == "n":
                        parts.append("\n"); i += 2; continue
                    elif nxt == "t":
                        parts.append("\t"); i += 2; continue
                    elif nxt == "r":
                        i += 2; continue
                    elif nxt in ("'", '"', "\\"):
                        parts.append(nxt); i += 2; continue
                parts.append(ch)
                i += 1
            return "".join(parts)

        for script in soup.find_all("script"):
            text = script.string or ""
            # Look for the content_noencode field inside cgiDataNew
            m = re.search(
                r"content_noencode\s*:\s*'((?:[^'\\]|\\.)*)'", text
            )
            if m:
                raw = m.group(1)
                decoded = _decode_js_escapes(raw)
                decoded = html_mod.unescape(decoded)
                return decoded

            # Alternative: the newer __APPMSG__ variable
            m = re.search(
                r"__APPMSG__\s*=\s*'((?:[^'\\]|\\.)*)'", text
            )
            if m:
                raw = m.group(1)
                decoded = _decode_js_escapes(raw)
                decoded = html_mod.unescape(decoded)
                import json as _json
                try:
                    data = _json.loads(decoded)
                    if isinstance(data, dict):
                        return data.get("content_noencode") or data.get("content") or decoded
                except Exception:
                    pass
                return decoded

        return None

    @staticmethod
    async def parse_url(url: str) -> str:
        """Fetch and parse a web page URL, extracting text content.

        Improved extraction that handles modern web pages (WeChat articles,
        news sites, SPAs) which often use <div>, <section>, <span> instead of
        semantic <p>/<article> tags.
        """
        import httpx
        from bs4 import BeautifulSoup, Tag

        # Full browser headers to avoid anti-bot detection
        chrome_ua = (
            "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 "
            "(KHTML, like Gecko) Chrome/131.0.0.0 Safari/537.36"
        )
        headers = {
            "User-Agent": chrome_ua,
            "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,"
                      "image/avif,image/webp,image/apng,*/*;q=0.8",
            "Accept-Language": "zh-CN,zh;q=0.9,en;q=0.8",
            "Accept-Encoding": "gzip, deflate, br",
            "Cache-Control": "no-cache",
            "Pragma": "no-cache",
            "Sec-Ch-Ua": '"Google Chrome";v="131", "Chromium";v="131", "Not_A Brand";v="24"',
            "Sec-Ch-Ua-Mobile": "?0",
            "Sec-Ch-Ua-Platform": '"Windows"',
            "Sec-Fetch-Dest": "document",
            "Sec-Fetch-Mode": "navigate",
            "Sec-Fetch-Site": "none" if "weixin" not in url else "cross-site",
            "Sec-Fetch-User": "?1",
            "Upgrade-Insecure-Requests": "1",
        }
        # WeChat-specific: add Referer to look like a normal click-through
        if "mp.weixin.qq.com" in url:
            headers["Referer"] = "https://weixin.qq.com/"

        async with httpx.AsyncClient(
            timeout=30.0,
            follow_redirects=True,
            headers=headers,
        ) as client:
            response = await client.get(url)
            response.raise_for_status()

            # Detect WeChat anti-bot verification page
            if "mp.weixin.qq.com" in url:
                text = response.text
                if "secitptpage/verify" in text or "环境异常" in text or "当前环境异常" in text:
                    # Retry once after a short delay (WeChat sometimes
                    # shows verify on first hit but passes on second)
                    import asyncio
                    logger.info("WeChat returned verify page, retrying after delay...")
                    await asyncio.sleep(2)
                    response = await client.get(url)
                    response.raise_for_status()
                    if "secitptpage/verify" in response.text or "环境异常" in response.text:
                        raise RuntimeError(
                            "WeChat requires verification (anti-bot). "
                            "Please wait a few minutes and try again, "
                            "or open the URL in a browser first to pass the check."
                        )

            soup = BeautifulSoup(response.text, "html.parser")

            # ── 1a. WeChat article? ──
            if "mp.weixin.qq.com" in url:
                wx_html = await DocumentParser._extract_wechat_article(soup)
                if wx_html:
                    # Re-parse the decoded article content
                    soup = BeautifulSoup(wx_html, "html.parser")
                    # Also try to extract meta from the original page
                    title_tag = soup.find("title") or soup.find("h1")
                    if not title_tag:
                        # Look for title in decoded content
                        for tag in soup.find_all(["h1", "h2", "h3"]):
                            if tag.get_text(strip=True):
                                title_tag = tag
                                break

            # ── 1. Remove noise elements ──
            for tag in soup([
                "script", "style", "nav", "footer", "header", "aside",
                "noscript", "iframe", "form", "button", "input",
            ]):
                tag.decompose()

            # Remove hidden elements
            for tag in soup.find_all(style=True):
                try:
                    style = (tag.get("style") or "").lower()
                except AttributeError:
                    continue
                if "display:none" in style or "display: none" in style or "visibility:hidden" in style:
                    tag.decompose()

            # Remove common nav/sidebar/footer classes
            noise_classes = [
                "nav", "navbar", "sidebar", "footer", "comment", "advertisement",
                "menu", "breadcrumb", "pagination", "related-posts", "recommend",
            ]
            for tag in soup.find_all(class_=True):
                try:
                    tag_classes = " ".join(tag.get("class") or []).lower()
                except AttributeError:
                    # Some tags (e.g. in WeChat pages) have attrs=None
                    continue
                if any(nc in tag_classes for nc in noise_classes):
                    try:
                        tag.decompose()
                    except Exception:
                        pass

            # ── 2. Try to find the main content area ──
            main_selectors = [
                # WeChat official accounts
                {"class_": "rich_media_content"},
                {"id": "js_content"},
                # General article selectors
                {"name": "article"},
                {"role": "main"},
                {"id": "content"},
                {"id": "article"},
                {"id": "main"},
                {"class_": "article"},
                {"class_": "post"},
                {"class_": "content"},
                {"class_": "article-content"},
                {"class_": "post-content"},
                {"class_": "entry-content"},
                {"class_": "markdown-body"},
                # Readme / docs
                {"class_": "readme"},
                {"class_": "documentation"},
            ]

            content_root = None
            for sel in main_selectors:
                content_root = soup.find(**sel)
                if content_root:
                    break

            if not content_root:
                content_root = soup.find("body") or soup

            # ── 3. Extract structured text ──
            # Tags to traverse in order of appearance.
            # We include div, section, span, article for modern pages
            # that don't use semantic tags.
            block_tags = {
                "p", "h1", "h2", "h3", "h4", "h5", "h6",
                "li", "pre", "blockquote", "td", "th", "figcaption",
            }
            inline_tags = {"div", "section", "span", "article", "details", "summary", "label"}

            paragraphs: list[str] = []
            seen_texts: set[str] = set()  # dedup identical blocks

            def _has_block_children(tag: Tag) -> bool:
                """Check if a tag or its descendants contain any block-level tags."""
                for child in tag.children:
                    if isinstance(child, Tag):
                        name = child.name.lower() if child.name else ""
                        if name in block_tags:
                            return True
                        if name in inline_tags and _has_block_children(child):
                            return True
                return False

            def _capture_text(tag: Tag) -> str | None:
                """Get all descendant text as a single string."""
                text = " ".join(
                    c.strip() for c in tag.strings
                    if isinstance(c, str) and c.strip()
                )
                return text if text and len(text) > 15 else None

            def extract_text(tag: Tag, depth: int = 0) -> None:
                """Recursively extract text blocks."""
                if depth > 12:
                    return

                tag_name = tag.name.lower() if tag.name else ""

                # Skip noise elements
                if tag_name in ("script", "style", "noscript", "iframe", "svg",
                                "img", "br", "hr", "button", "input", "form"):
                    return

                # Block-level elements: extract their text as a unit
                if tag_name in block_tags:
                    text = tag.get_text(" ", strip=True)
                    if text and len(text) > 5 and text not in seen_texts:
                        seen_texts.add(text)
                        if tag_name.startswith("h"):
                            level = int(tag_name[1])
                            paragraphs.append(f"{'#' * level} {text}")
                        elif tag_name == "li":
                            paragraphs.append(f"- {text}")
                        else:
                            paragraphs.append(text)
                    return

                # Inline container tags: recurse into children first.
                # Only capture as a leaf if there are no useful children.
                if tag_name in inline_tags:
                    # Collect child tags (not NavigableString)
                    child_tags = [c for c in tag.children if isinstance(c, Tag)]
                    child_tags = [c for c in child_tags
                                  if c.name and c.name.lower() not in
                                  ("script", "style", "noscript", "iframe", "svg",
                                   "img", "br", "hr", "button", "input")]

                    if child_tags:
                        # Has structured children — recurse into them
                        for child in child_tags:
                            extract_text(child, depth + 1)
                    else:
                        # Leaf container — capture all descendant text
                        text = _capture_text(tag)
                        if text and text not in seen_texts:
                            seen_texts.add(text)
                            paragraphs.append(text)
                    return

                # Unknown tags — recurse
                for child in tag.children:
                    if isinstance(child, Tag):
                        extract_text(child, depth + 1)

            extract_text(content_root)

            # ── 4. Fallback if nothing extracted ──
            if not paragraphs:
                text = content_root.get_text(separator="\n", strip=True)
                # Clean up excessive whitespace
                import re
                text = re.sub(r"\n{3,}", "\n\n", text)
                return text

            return "\n\n".join(paragraphs)

    @staticmethod
    def get_file_type(filename: str) -> str:
        """Get normalized file type from filename."""
        ext = Path(filename).suffix.lower().lstrip(".")
        type_map = {
            "pdf": "pdf",
            "docx": "docx", "doc": "docx",
            "xlsx": "xlsx", "xls": "xlsx",
            "csv": "csv",
            "txt": "txt", "md": "md", "markdown": "md",
            "pptx": "pptx", "ppt": "pptx",
            "png": "image", "jpg": "image", "jpeg": "image",
            "gif": "image", "webp": "image", "bmp": "image", "svg": "image",
            "mp3": "audio", "wav": "audio", "m4a": "audio", "aac": "audio",
            "ogg": "audio", "flac": "audio",
            "mp4": "video", "mov": "video", "avi": "video", "mkv": "video", "webm": "video",
            "html": "html", "htm": "html",
            "xml": "txt", "json": "txt", "log": "txt",
            "properties": "txt", "yaml": "txt", "yml": "txt", "toml": "txt",
            "ini": "txt", "cfg": "txt",
        }
        return type_map.get(ext, "txt")
